"""多格式文档加载器 — 多引擎兜底（pypdf / docutils / unstructured / MarkItDown）."""

import csv
import io
import json
import re
from pathlib import Path

from src.config import settings
from src.utils.logger import logger

TEXT_EXTENSIONS = {".txt", ".md", ".markdown", ".rst"}
TEXT_ENCODINGS = ("utf-8", "utf-8-sig", "gb18030", "gbk")


def _read_text_with_fallback(path: Path) -> str:
    for encoding in TEXT_ENCODINGS:
        try:
            return path.read_text(encoding=encoding)
        except UnicodeDecodeError:
            continue
    logger.warning("text_decode_replacement_fallback", file=path.name)
    return path.read_text(encoding="utf-8", errors="replace")


SUPPORTED_EXTENSIONS = {
    ".pdf", ".docx", ".doc", ".md", ".markdown",
    ".html", ".htm", ".txt", ".rst",
    ".csv", ".xlsx", ".xls", ".json", ".pptx", ".epub",
}


def load_document(file_path: str) -> list[dict]:
    path = Path(file_path)
    ext = path.suffix.lower()

    if ext not in SUPPORTED_EXTENSIONS:
        raise ValueError(f"Unsupported file type '{ext}'. Supported: {sorted(SUPPORTED_EXTENSIONS)}")

    # ── 纯文本 ────────────────────────────────────
    if ext == ".txt":
        text = _read_text_with_fallback(path)
        return [{"content": text, "metadata": {"source": path.name, "page": 1}}]

    # ── Markdown ──────────────────────────────────
    if ext in (".md", ".markdown"):
        text = _read_text_with_fallback(path)
        sections = _split_markdown_sections(text)
        return [{"content": s["content"], "metadata": {"source": path.name, "section": s.get("heading", ""), "page": 1}} for s in sections]

    # ── PDF — pypdf 引擎，零额外依赖 ─────────────
    if ext == ".pdf":
        return _load_pdf(path)

    # ── RST — docutils 引擎 ───────────────────────
    if ext == ".rst":
        return _load_rst(path)

    # ── CSV ───────────────────────────────────────
    if ext == ".csv":
        return _load_csv(path)

    # ── Excel (.xlsx / .xls) ──────────────────────
    if ext in (".xlsx", ".xls"):
        return _load_xlsx(path)

    # ── JSON ──────────────────────────────────────
    if ext == ".json":
        return _load_json(path)

    # ── PPTX ──────────────────────────────────────
    if ext == ".pptx":
        return _load_pptx(path)

    # ── EPUB ──────────────────────────────────────
    if ext == ".epub":
        return _load_epub(path)

    # ── DOCX / DOC / HTML — unstructured 引擎 ─────
    try:
        from unstructured.partition.auto import partition
    except ImportError:
        raise ImportError(
            "unstructured library required for non-text files. Install: pip install unstructured"
        )

    elements = partition(filename=file_path)
    docs = []
    for el in elements:
        docs.append({
            "content": str(el),
            "metadata": {
                "source": path.name,
                "page": getattr(el.metadata, "page_number", 1) if el.metadata else 1,
            },
        })
    logger.info("document_loaded", file=path.name, elements=len(docs))
    return docs


# ═══════════════════════════════════════════════════
#  PDF 解析 — MarkItDown 主引擎 + pypdf 兜底
# ═══════════════════════════════════════════════════

def _load_pdf(path: Path) -> list[dict]:
    # 读取配置决定引擎
    engine = settings.pdf_engine
    if engine == "pypdf":
        return _load_pdf_pypdf(path)

    # 默认 markitdown — 尝试导入
    try:
        from markitdown import MarkItDown
    except ImportError:
        logger.warning("markitdown_not_installed_fallback_pypdf", file=path.name)
        return _load_pdf_pypdf(path)

    try:
        md = MarkItDown()
        result = md.convert(str(path))
        text = result.text_content
    except Exception as e:
        logger.warning("markitdown_failed_fallback_pypdf", file=path.name, error=str(e))
        return _load_pdf_pypdf(path)

    if not text or not text.strip():
        logger.warning("pdf_empty", file=path.name)
        return [{"content": "", "metadata": {"source": path.name, "page": 1}}]

    # 通过 Markdown 章节切分后处理
    sections = _split_markdown_sections(text)
    docs = [
        {"content": s["content"], "metadata": {"source": path.name, "section": s.get("heading", ""), "page": 1}}
        for s in sections
    ]
    logger.info("document_loaded", file=path.name, elements=len(docs), engine="markitdown")
    return docs


def _load_pdf_pypdf(path: Path) -> list[dict]:
    try:
        from pypdf import PdfReader
    except ImportError:
        raise ImportError("pypdf required for PDF files. Install: pip install pypdf")

    reader = PdfReader(str(path))
    docs = []
    for i, page in enumerate(reader.pages, start=1):
        text = page.extract_text()
        if text and text.strip():
            docs.append({
                "content": text.strip(),
                "metadata": {"source": path.name, "page": i},
            })
    if not docs:
        logger.warning("pdf_empty", file=path.name)
        return [{"content": "", "metadata": {"source": path.name, "page": 1}}]
    logger.info("document_loaded", file=path.name, elements=len(docs), engine="pypdf")
    return docs


# ═══════════════════════════════════════════════════
#  RST 解析 — docutils
# ═══════════════════════════════════════════════════

def _load_rst(path: Path) -> list[dict]:
    try:
        from docutils.core import publish_parts
    except ImportError:
        raise ImportError("docutils required for RST files. Install: pip install docutils")

    rst_text = _read_text_with_fallback(path)
    try:
        html = publish_parts(rst_text, writer_name="html")["body"]
        text = _strip_html(html)
    except Exception:
        # 降级：直接当纯文本读取
        text = rst_text
        logger.warning("rst_docutils_failed_plain_text", file=path.name)

    if not text.strip():
        return [{"content": rst_text, "metadata": {"source": path.name, "page": 1}}]

    logger.info("document_loaded", file=path.name, engine="docutils")
    return [{"content": text.strip(), "metadata": {"source": path.name, "page": 1}}]


def _strip_html(html: str) -> str:
    """简单去除 HTML 标签，保留文本内容。"""
    text = re.sub(r"<[^>]+>", " ", html)
    text = re.sub(r"&[a-z]+;", " ", text)
    text = re.sub(r"\s+", " ", text)
    return text


def _split_markdown_sections(text: str) -> list[dict]:
    lines = text.split("\n")
    sections = []
    current_heading = ""
    current_lines = []

    for line in lines:
        if line.startswith("#"):
            if current_lines:
                sections.append({"heading": current_heading, "content": "\n".join(current_lines).strip()})
            current_heading = line.lstrip("#").strip()
            current_lines = []
        else:
            current_lines.append(line)

    if current_lines:
        sections.append({"heading": current_heading, "content": "\n".join(current_lines).strip()})

    if not sections:
        sections.append({"heading": "", "content": text})

    return sections


# ═══════════════════════════════════════════════════
#  CSV 加载
# ═══════════════════════════════════════════════════

def _load_csv(path: Path | str) -> list[dict]:
    path = Path(path)
    raw_text = _read_text_with_fallback(path)
    # 尝试自动检测 dialect
    try:
        dialect = csv.Sniffer().sniff(raw_text[:4096])
        has_header = csv.Sniffer().has_header(raw_text[:4096])
    except csv.Error:
        dialect = "excel"
        has_header = False

    reader = csv.reader(io.StringIO(raw_text), dialect)
    rows = list(reader)
    if not rows:
        logger.warning("csv_empty", file=path.name)
        return [{"content": "", "metadata": {"source": path.name}}]

    docs = []
    if has_header:
        headers = rows[0]
        data_rows = rows[1:]
    else:
        headers = [f"col_{i}" for i in range(len(rows[0]))]
        data_rows = rows

    for ri, row in enumerate(data_rows, start=1):
        parts = []
        for ci, val in enumerate(row):
            col_name = headers[ci] if ci < len(headers) else f"col_{ci}"
            parts.append(f"{col_name}: {val}")
        content = ", ".join(parts)
        if content.strip():
            docs.append({
                "content": content,
                "metadata": {"source": path.name, "row": ri, "columns": headers},
            })

    logger.info("document_loaded", file=path.name, elements=len(docs), engine="csv")
    return docs


# ═══════════════════════════════════════════════════
#  Excel (.xlsx / .xls) 加载
# ═══════════════════════════════════════════════════

def _load_xlsx(path: Path | str) -> list[dict]:
    path = Path(path)
    try:
        import openpyxl
    except ImportError:
        raise ImportError("openpyxl required for Excel files. Install: pip install openpyxl")

    wb = openpyxl.load_workbook(str(path), data_only=True)
    docs = []

    for si, sheet_name in enumerate(wb.sheetnames):
        ws = wb[sheet_name]
        # 处理合并单元格：将合并区域的值传播到所有单元格
        _propagate_merged_cells(ws)

        rows = list(ws.iter_rows(values_only=True))
        if not rows:
            continue

        headers = rows[0] if rows else []
        headers = [str(h) if h is not None else f"col_{i}" for i, h in enumerate(headers)]
        data_rows = rows[1:]

        for ri, row in enumerate(data_rows, start=2):
            parts = []
            for ci, val in enumerate(row):
                col_name = headers[ci] if ci < len(headers) else f"col_{ci}"
                parts.append(f"{col_name}: {val}")
            content = ", ".join(parts)
            if content.strip():
                docs.append({
                    "content": content,
                    "metadata": {
                        "source": path.name,
                        "row": ri,
                        "sheet_name": sheet_name,
                        "sheet_index": si,
                        "columns": headers,
                    },
                })

    if not docs:
        logger.warning("xlsx_empty", file=path.name)
        return [{"content": "", "metadata": {"source": path.name}}]
    logger.info("document_loaded", file=path.name, elements=len(docs), engine="openpyxl")
    return docs


def _propagate_merged_cells(ws) -> None:
    """将 openpyxl 合并单元格的值传播到所有受影响的位置。"""
    # 收集所有合并范围
    merged_ranges = list(ws.merged_cells.ranges)
    for merged_range in merged_ranges:
        min_col = merged_range.min_col
        max_col = merged_range.max_col
        min_row = merged_range.min_row
        max_row = merged_range.max_row
        # 获取左上角单元格的值
        top_left_val = ws.cell(row=min_row, column=min_col).value
        # 传播到整个合并区域
        for r in range(min_row, max_row + 1):
            for c in range(min_col, max_col + 1):
                if ws.cell(row=r, column=c).value is None:
                    ws.cell(row=r, column=c).value = top_left_val


# ═══════════════════════════════════════════════════
#  JSON 加载
# ═══════════════════════════════════════════════════

def _load_json(path: Path | str) -> list[dict]:
    path = Path(path)
    raw_text = _read_text_with_fallback(path)
    try:
        data = json.loads(raw_text)
    except json.JSONDecodeError as e:
        raise ValueError(f"Invalid JSON file '{path.name}': {e}")

    docs = []
    # 如果是对象数组，每个元素一个 chunk
    if isinstance(data, list):
        for i, item in enumerate(data):
            content = _flatten_json_to_text(item)
            if content.strip():
                docs.append({
                    "content": content,
                    "metadata": {"source": path.name, "index": i, "json_path": f"$[{i}]"},
                })
    else:
        # 单个对象，包装为单 chunk
        content = _flatten_json_to_text(data)
        docs.append({
            "content": content,
            "metadata": {"source": path.name, "index": 0, "json_path": "$"},
        })

    if not docs:
        logger.warning("json_empty", file=path.name)
        return [{"content": "", "metadata": {"source": path.name}}]
    logger.info("document_loaded", file=path.name, elements=len(docs), engine="json")
    return docs


def _flatten_json_to_text(obj, prefix: str = "", max_depth: int = 10) -> str:
    """递归展平 JSON 对象为 key: value 文本，使用 dot-notation 键名。"""
    if max_depth <= 0:
        return str(obj)

    if isinstance(obj, dict):
        parts = []
        for k, v in obj.items():
            key = f"{prefix}.{k}" if prefix else k
            if isinstance(v, (dict, list)):
                parts.append(_flatten_json_to_text(v, key, max_depth - 1))
            else:
                parts.append(f"{key}: {v}")
        return "; ".join(parts)
    elif isinstance(obj, list):
        # 对于原始值列表
        if all(not isinstance(x, (dict, list)) for x in obj):
            return f"{prefix}: [{', '.join(str(x) for x in obj)}]"
        # 对于对象列表，逐个处理
        parts = []
        for i, item in enumerate(obj):
            key = f"{prefix}[{i}]" if prefix else f"[{i}]"
            parts.append(_flatten_json_to_text(item, key, max_depth - 1))
        return "; ".join(parts)
    else:
        return f"{prefix}: {obj}"


# ═══════════════════════════════════════════════════
#  PPTX 加载
# ═══════════════════════════════════════════════════

def _load_pptx(path: Path | str) -> list[dict]:
    path = Path(path)
    try:
        from pptx import Presentation
    except ImportError:
        raise ImportError("python-pptx required for PPTX files. Install: pip install python-pptx")

    prs = Presentation(str(path))
    docs = []

    for si, slide in enumerate(prs.slides, start=1):
        title_text = ""
        paragraphs = []

        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    para_text = para.text.strip()
                    if not para_text:
                        continue
                    # 标题占位符优先
                    if shape.is_placeholder and shape.placeholder_format.type == 1:  # TITLE
                        title_text = para_text
                    else:
                        paragraphs.append(para_text)

        if not title_text and paragraphs:
            title_text = paragraphs[0]
            paragraphs = paragraphs[1:]

        content = f"标题: {title_text}\n" + "\n".join(paragraphs) if title_text else "\n".join(paragraphs)
        content = content.strip()
        if content:
            docs.append({
                "content": content,
                "metadata": {"source": path.name, "slide": si},
            })

    if not docs:
        logger.warning("pptx_empty", file=path.name)
        return [{"content": "", "metadata": {"source": path.name}}]
    logger.info("document_loaded", file=path.name, elements=len(docs), engine="python-pptx")
    return docs


# ═══════════════════════════════════════════════════
#  EPUB 加载
# ═══════════════════════════════════════════════════

def _load_epub(path: Path | str) -> list[dict]:
    path = Path(path)
    try:
        from ebooklib import epub
        from bs4 import BeautifulSoup
    except ImportError:
        raise ImportError("ebooklib and beautifulsoup4 required for EPUB files. Install: pip install ebooklib beautifulsoup4")

    book = epub.read_epub(str(path))
    docs = []

    for item in book.get_items():
        if item.get_type() == 9:  # ITEM_DOCUMENT (HTML)
            soup = BeautifulSoup(item.get_content(), "html.parser")
            # 提取纯文本
            text = soup.get_text(separator="\n")
            # 清理多余空行
            lines = [line.strip() for line in text.split("\n") if line.strip()]
            clean_text = "\n".join(lines)
            if not clean_text.strip():
                continue

            # 按标题拆分章节
            sections = _split_markdown_sections(clean_text)
            for s in sections:
                item_name = item.get_name() or "unknown"
                docs.append({
                    "content": s["content"],
                    "metadata": {
                        "source": path.name,
                        "chapter": item_name,
                        "section": s.get("heading", ""),
                    },
                })

    if not docs:
        # 单文档 EPUB — 合并所有文本
        all_text = []
        for item in book.get_items():
            if item.get_type() == 9:
                soup = BeautifulSoup(item.get_content(), "html.parser")
                text = soup.get_text(separator="\n")
                all_text.append(text)
        full_text = "\n".join(all_text)
        if full_text.strip():
            sections = _split_markdown_sections(full_text)
            docs = [
                {"content": s["content"], "metadata": {"source": path.name, "section": s.get("heading", "")}}
                for s in sections
            ]

    if not docs:
        logger.warning("epub_empty", file=path.name)
        return [{"content": "", "metadata": {"source": path.name}}]
    logger.info("document_loaded", file=path.name, elements=len(docs), engine="ebooklib")
    return docs
